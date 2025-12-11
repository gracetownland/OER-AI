"""
Media Processing Job for OER Textbook Media Items
Processes media items (video transcripts, PDFs, PPTs) attached to chapters
"""

import requests
import boto3
import json
import psycopg2
import sys
import logging
from datetime import datetime
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_aws import BedrockEmbeddings
from langchain_postgres import PGVector
from langchain_core.documents import Document
from awsglue.utils import getResolvedOptions
from awsglue.context import GlueContext
from pyspark.context import SparkContext
from urllib.parse import urlparse
import time
import io
from typing import List, Dict, Optional
from bs4 import BeautifulSoup

# Global variables
connection = None
db_secret = None
embeddings = None
vector_store = None

# Database configuration
DB_SECRET_NAME = None
RDS_PROXY_ENDPOINT = None
EMBEDDING_MODEL_ID = None

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

print("=== MEDIA PROCESSING JOB START ===")

# Get job parameters
try:
    args = getResolvedOptions(sys.argv, [
        'batch_id',
        'sqs_message_id',
        'sqs_message_body',
        'trigger_timestamp',
        'region_name',
        'GLUE_BUCKET',
        'rds_secret',
        'rds_proxy_endpoint',
        'embedding_model_id',
        'media_url',
        'media_type'
    ])
    sc = SparkContext()
    glueContext = GlueContext(sc)
    print("=== JOB PARAMETERS ===")
    for key, value in args.items():
        print(f"{key}: {value}")
        
    # Initialize database configuration
    DB_SECRET_NAME = args['rds_secret']
    RDS_PROXY_ENDPOINT = args['rds_proxy_endpoint']
    EMBEDDING_MODEL_ID = args['embedding_model_id']
        
    # Parse the SQS message body
    sqs_data = json.loads(args['sqs_message_body'])
    metadata = sqs_data.get('metadata', {})
    
    # Extract fields from message
    media_url = args.get('media_url') or sqs_data.get('media_url')
    media_type = args.get('media_type') or sqs_data.get('media_type')
    
    # Extract metadata from CSV upload
    book_title = metadata.get('book_title', '')
    media_title = metadata.get('media_title', '')
    chapter_title = metadata.get('chapter_title', '')
    chapter_url = metadata.get('chapter_url', '')
    media_type_raw = metadata.get('media_type_raw', '')
    
    print(f"=== PROCESSING MEDIA ITEM ===")
    print(f"Book Title: {book_title}")
    print(f"Media Title: {media_title}")
    print(f"Media Type: {media_type}")
    print(f"Media Type Raw: {media_type_raw}")
    print(f"Chapter Title: {chapter_title}")
    print(f"Chapter URL: {chapter_url}")
    print(f"Media URL: {media_url}")
    
except Exception as e:
    print(f"Error parsing arguments: {e}")
    sys.exit(1)


# Initialize AWS clients
secrets_manager = boto3.client("secretsmanager", region_name=args['region_name'])
s3_client = boto3.client("s3", region_name=args['region_name'])

def get_secret(secret_name, expect_json=True):
    global db_secret
    if db_secret is None:
        try:
            response = secrets_manager.get_secret_value(SecretId=secret_name)["SecretString"]
            db_secret = json.loads(response) if expect_json else response
        except Exception as e:
            logger.error(f"Error fetching secret: {e}")
            raise
    return db_secret

def connect_to_db():
    global connection
    if connection is None or connection.closed:
        try:
            secret = get_secret(DB_SECRET_NAME)
            connection = psycopg2.connect(
                dbname=secret["dbname"],
                user=secret["username"],
                password=secret["password"],
                host=RDS_PROXY_ENDPOINT,
                port=int(secret["port"])
            )
            logger.info("Connected to the database!")
        except Exception as e:
            logger.error(f"Failed to connect to database: {e}")
            raise
    return connection

def initialize_embeddings_and_vectorstore(textbook_id: str, textbook_title: str):
    """
    Initialize Bedrock embeddings and PGVector store for the textbook.
    Returns the initialized vector store.
    """
    global embeddings, vector_store
    
    try:
        # Initialize Bedrock embeddings if not already done
        if embeddings is None:
            logger.info(f"Initializing Bedrock embeddings with model: {EMBEDDING_MODEL_ID}")
            embeddings = BedrockEmbeddings(
                model_id=EMBEDDING_MODEL_ID,
                region_name='us-east-1',
                model_kwargs={"input_type": "search_document"}
            )
            logger.info("Bedrock embeddings initialized successfully")
        
        # Get database connection info for vector store
        conn_info = get_secret(DB_SECRET_NAME)
        conn_str = (
            f"postgresql://{conn_info['username']}:{conn_info['password']}"
            f"@{RDS_PROXY_ENDPOINT}:{conn_info['port']}/{conn_info['dbname']}"
        )
        
        # Initialize PGVector store
        logger.info(f"Initializing PGVector store for textbook ID: {textbook_id}")
        vector_store = PGVector(
            embeddings=embeddings,
            collection_name=str(textbook_id),
            collection_metadata={'title': textbook_title},
            connection=conn_str,
            use_jsonb=True
        )
        logger.info("PGVector store initialized successfully")
        
        return vector_store
        
    except Exception as e:
        logger.error(f"Error initializing embeddings and vector store: {e}")
        raise

def execute_query(query, params=None, fetch_one=False):
    """Execute a database query and return results"""
    conn = None
    cursor = None
    try:
        conn = connect_to_db()
        cursor = conn.cursor()
        cursor.execute(query, params)
        
        if fetch_one:
            result = cursor.fetchone()
        else:
            result = cursor.fetchall() if cursor.description else None
            
        conn.commit()
        return result
    except Exception as e:
        if conn:
            conn.rollback()
        logger.error(f"Query execution error: {e}")
        raise
    finally:
        if cursor:
            cursor.close()

def get_textbook_by_id(textbook_id: str) -> Optional[Dict]:
    """Get textbook information by ID"""
    query = """
        SELECT id, title, source_url, metadata
        FROM textbooks
        WHERE id = %s
    """
    result = execute_query(query, (textbook_id,), fetch_one=True)
    if result:
        return {
            'id': str(result[0]),
            'title': result[1],
            'source_url': result[2],
            'metadata': result[3]
        }
    return None

def get_section_by_url(source_url: str) -> Optional[Dict]:
    """Get section information by source URL"""
    query = """
        SELECT id, textbook_id, title, order_index, source_url
        FROM sections
        WHERE source_url = %s
        LIMIT 1
    """
    result = execute_query(query, (source_url,), fetch_one=True)
    if result:
        return {
            'id': str(result[0]),
            'textbook_id': str(result[1]),
            'title': result[2],
            'order_index': result[3],
            'source_url': result[4]
        }
    return None

def map_media_type_to_db_enum(media_type: str) -> str:
    """
    Map processing media types to database enum values.
    
    Database enum: ('pdf', 'audio', 'video', 'image', 'transcript', 'h5p', 'other')
    Processing types: 'pdf', 'pptx', 'ppt', 'video_transcript'
    """
    type_mapping = {
        'pdf': 'pdf',
        'pptx': 'other',  # PowerPoint files stored as 'other'
        'ppt': 'other',   # PowerPoint files stored as 'other'
        'video_transcript': 'transcript',  # H5P transcripts
        'video': 'video',
        'audio': 'audio',
        'image': 'image',
        'h5p': 'h5p'
    }
    
    db_type = type_mapping.get(media_type, 'other')
    
    if db_type != media_type:
        logger.info(f"Mapped media_type '{media_type}' to database enum '{db_type}'")
    
    return db_type

def create_or_update_media_item(textbook_id: str, section_id: str, media_type: str, 
                                 media_url: str, source_url: str, description: str = None) -> str:
    """
    Create or update a media item in the database.
    Returns the media_item_id.
    """
    try:
        # Map media_type to database enum value
        db_media_type = map_media_type_to_db_enum(media_type)
        
        # Check if media item already exists
        check_query = """
            SELECT id FROM media_items
            WHERE textbook_id = %s AND section_id = %s AND source_url = %s
        """
        existing = execute_query(check_query, (textbook_id, section_id, source_url), fetch_one=True)
        
        if existing:
            # Update existing media item
            media_item_id = str(existing[0])
            update_query = """
                UPDATE media_items
                SET media_type = %s, uri = %s, description = %s
                WHERE id = %s
            """
            execute_query(update_query, (db_media_type, media_url, description, media_item_id))
            logger.info(f"Updated existing media item: {media_item_id}")
        else:
            # Create new media item
            insert_query = """
                INSERT INTO media_items (textbook_id, section_id, media_type, uri, source_url, description)
                VALUES (%s, %s, %s, %s, %s, %s)
                RETURNING id
            """
            result = execute_query(
                insert_query,
                (textbook_id, section_id, db_media_type, media_url, source_url, description),
                fetch_one=True
            )
            media_item_id = str(result[0])
            logger.info(f"Created new media item: {media_item_id}")
        
        return media_item_id
        
    except Exception as e:
        logger.error(f"Error creating/updating media item: {e}")
        raise

def scrape_transcript_url(page_url: str) -> Optional[str]:
    """
    Scrape the page to find the transcript download link.
    Uses CSS selector: #attachments-tab > table > tbody > tr:nth-child(3) > td:nth-child(6) > div > a
    """
    try:
        logger.info(f"Scraping transcript URL from: {page_url}")
        response = requests.get(page_url, timeout=30)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Try the specific CSS selector first
        selector = "#attachments-tab > table > tbody > tr:nth-child(3) > td:nth-child(6) > div > a"
        link_element = soup.select_one(selector)
        
        if link_element and link_element.get('href'):
            transcript_url = link_element['href']
            # Make absolute URL if relative
            if not transcript_url.startswith('http'):
                from urllib.parse import urljoin
                transcript_url = urljoin(page_url, transcript_url)
            logger.info(f"Found transcript URL: {transcript_url}")
            return transcript_url
        
        # Fallback: look for .txt file links in the page
        logger.warning("Specific selector not found, searching for .txt links")
        for link in soup.find_all('a', href=True):
            href = link['href']
            if href.endswith('.txt'):
                if not href.startswith('http'):
                    from urllib.parse import urljoin
                    href = urljoin(page_url, href)
                logger.info(f"Found .txt file: {href}")
                return href
        
        logger.error("No transcript URL found on page")
        return None
        
    except Exception as e:
        logger.error(f"Error scraping transcript URL: {e}")
        return None

def download_file_from_url(url: str) -> bytes:
    """Download a file from a URL"""
    try:
        logger.info(f"Downloading file from: {url}")
        response = requests.get(url, timeout=60)
        response.raise_for_status()
        logger.info(f"Downloaded {len(response.content)} bytes")
        return response.content
    except Exception as e:
        logger.error(f"Error downloading file from {url}: {e}")
        raise

def download_file_from_s3(s3_uri: str) -> bytes:
    """Download a file from S3"""
    try:
        # Parse S3 URI (s3://bucket/key)
        parsed = urlparse(s3_uri)
        bucket = parsed.netloc
        key = parsed.path.lstrip('/')
        
        logger.info(f"Downloading from S3: bucket={bucket}, key={key}")
        response = s3_client.get_object(Bucket=bucket, Key=key)
        content = response['Body'].read()
        logger.info(f"Downloaded {len(content)} bytes from S3")
        return content
    except Exception as e:
        logger.error(f"Error downloading file from S3 {s3_uri}: {e}")
        raise

def extract_text_from_transcript(content: bytes) -> str:
    """Extract text from transcript file (plain text)"""
    try:
        text = content.decode('utf-8')
        logger.info(f"Extracted {len(text)} characters from transcript")
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from transcript: {e}")
        raise

def extract_text_from_pdf(content: bytes) -> str:
    """Extract text from PDF file"""
    try:
        from io import BytesIO
        import PyPDF2
        
        pdf_file = BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page_num, page in enumerate(pdf_reader.pages):
            text += page.extract_text() + "\n\n"
        
        logger.info(f"Extracted {len(text)} characters from PDF ({len(pdf_reader.pages)} pages)")
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        raise

def extract_text_from_pptx(content: bytes) -> str:
    """Extract text from PowerPoint file"""
    try:
        from io import BytesIO
        from pptx import Presentation
        
        pptx_file = BytesIO(content)
        prs = Presentation(pptx_file)
        text = ""
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        
        logger.info(f"Extracted {len(text)} characters from PPTX ({len(prs.slides)} slides)")
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        raise

def process_media_item(book_title: str, media_title: str, chapter_title: str, 
                       chapter_url: str, media_type: str, media_url: str) -> int:
    """
    Process a media item and store its chunks and embeddings.
    Returns the number of chunks created.
    """
    try:
        # Step 1: Find section (which contains textbook_id)
        logger.info(f"Looking up section by URL: {chapter_url}")
        section_info = get_section_by_url(chapter_url)
        if not section_info:
            logger.error(f"Section with URL '{chapter_url}' not found in database")
            return 0
        
        section_id = section_info['id']
        textbook_id = section_info['textbook_id']
        logger.info(f"Found section ID: {section_id}")
        logger.info(f"Textbook ID from section: {textbook_id}")
        
        # Step 2: Get textbook information
        logger.info(f"Looking up textbook by ID: {textbook_id}")
        textbook_info = get_textbook_by_id(textbook_id)
        if not textbook_info:
            logger.error(f"Textbook ID '{textbook_id}' not found in database")
            return 0
        
        # Step 3: Create or update media item in database
        logger.info(f"Creating/updating media item in database")
        media_item_id = create_or_update_media_item(
            textbook_id=textbook_id,
            section_id=section_id,
            media_type=media_type,
            media_url=media_url,
            source_url=media_url,  # Use media_url as source_url
            description=media_title
        )
        
        # Step 4: Download and extract text from file
        logger.info(f"Processing media type: {media_type}")
        
        # TODO: H5P video transcripts require Selenium/browser automation
        # For now, skip processing but don't fail
        if media_type == 'video_transcript':
            logger.warning("H5P video transcript processing not yet implemented")
            logger.warning("Skipping text extraction for this media item")
            logger.info(f"Media item created in database with ID: {media_item_id}")
            logger.info("To enable H5P processing, implement transcript resolver Lambda with Selenium")
            return 0  # Return 0 chunks created, but don't raise error
        elif media_url.startswith('s3://'):
            content = download_file_from_s3(media_url)
            # Extract text based on media type
            if media_type == 'pdf':
                text = extract_text_from_pdf(content)
            elif media_type in ['pptx', 'ppt']:
                text = extract_text_from_pptx(content)
            else:
                logger.warning(f"Unsupported media type for S3: {media_type}")
                return 0
        else:
            # HTTP/HTTPS URL - Direct download link
            # For PDFs: Opens in PDF viewer (direct link to file)
            # For PPTX: Direct download link to PowerPoint file
            content = download_file_from_url(media_url)
            # Extract text based on media type
            if media_type == 'pdf':
                text = extract_text_from_pdf(content)
            elif media_type in ['pptx', 'ppt']:
                text = extract_text_from_pptx(content)
            else:
                logger.warning(f"Unsupported media type: {media_type}")
                return 0
        
        
        if not text.strip():
            logger.warning(f"No text extracted from media item")
            return 0
        
        logger.info(f"Extracted {len(text)} characters from media item")
        
        # Step 5: Initialize embeddings and vector store
        initialize_embeddings_and_vectorstore(textbook_id, textbook_info['title'])
        
        # Step 6: Split text into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        
        chunks = text_splitter.split_text(text)
        logger.info(f"Split text into {len(chunks)} chunks")
        
        # Step 7: Create document chunks with metadata
        documents = []
        for i, chunk_text in enumerate(chunks):
            metadata = {
                'textbook_id': textbook_id,
                'textbook_title': textbook_info['title'],
                'section_id': section_id,
                'section_title': section_info['title'],
                'section_order': section_info['order_index'],
                'media_item_id': media_item_id,
                'media_type': media_type,
                'media_title': media_title,
                'source': media_url,  # Source URL of the media file
                'chunk_index': i,
                'total_chunks': len(chunks)
            }
            
            documents.append(Document(
                page_content=chunk_text,
                metadata=metadata
            ))
        
        # Step 8: Store chunks in database
        chunk_ids = []
        for doc in documents:
            query = """
                INSERT INTO document_chunks (textbook_id, section_id, media_item_id, chunk_text, chunk_meta)
                VALUES (%s, %s, %s, %s, %s)
                RETURNING id
            """
            result = execute_query(
                query,
                (
                    textbook_id,
                    section_id,
                    media_item_id,
                    doc.page_content,
                    json.dumps(doc.metadata)
                ),
                fetch_one=True
            )
            if result:
                chunk_ids.append(str(result[0]))
        
        logger.info(f"Created {len(chunk_ids)} chunks in database")
        
        # Step 9: Add documents to vector store
        if vector_store and documents:
            try:
                vector_store.add_documents(documents)
                logger.info(f"Added {len(documents)} documents to vector store")
            except Exception as e:
                logger.error(f"Error adding documents to vector store: {e}")
                # Continue even if vector store fails
        
        return len(chunks)
        
    except Exception as e:
        logger.error(f"Error processing media item: {e}")
        raise

# Main execution
try:
    logger.info("Starting media processing...")
    
    chunks_created = process_media_item(
        book_title=book_title,
        media_title=media_title,
        chapter_title=chapter_title,
        chapter_url=chapter_url,
        media_type=media_type,
        media_url=media_url
    )
    
    logger.info(f"Media processing completed successfully. Created {chunks_created} chunks.")
    
except Exception as e:
    logger.error(f"Media processing failed: {e}")
    import traceback
    traceback.print_exc()
    sys.exit(1)
finally:
    if connection and not connection.closed:
        connection.close()
        logger.info("Database connection closed")
    if sc:
        sc.stop()
        logger.info("Spark context stopped")

print("=== MEDIA PROCESSING JOB END ===")
